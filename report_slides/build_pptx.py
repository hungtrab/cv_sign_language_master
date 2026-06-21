"""Build the English-language ASL pipeline capstone deck as a .pptx.

Mirrors the content/structure of Slide_ASL.tex (HUST beamer deck), translated
to English, rebuilt natively in PowerPoint with a simple navy/blue academic
style (no decorative icons, consistent margins, left-aligned body text).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import glob
import os

NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x1E, 0x7B, 0x34)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
HEADER_H = Inches(0.9)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

page_num = 0


def new_slide():
    return prs.slides.add_slide(BLANK)


def add_header(slide, title_text, section=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = MARGIN
    tf.margin_top = Pt(6)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    # thin accent line under header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, HEADER_H, SLIDE_W, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    line.shadow.inherit = False
    if section:
        sec = slide.shapes.add_textbox(SLIDE_W - Inches(3.2), Pt(10), Inches(2.8), Inches(0.4))
        p2 = sec.text_frame.paragraphs[0]
        p2.text = section
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p2.alignment = PP_ALIGN.RIGHT


def add_footer(slide):
    global page_num
    page_num += 1
    box = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.45), Inches(0.7), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items, left=MARGIN, top=Inches(1.15), width=None, height=None,
                 font_size=18, bold_first=False):
    """items: list of (text, level, color_or_None, bold_or_None) or plain strings."""
    width = width or (SLIDE_W - 2 * MARGIN)
    height = height or (SLIDE_H - top - Inches(0.6))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, level, color, bold = item, 0, None, False
        else:
            text, level, color, bold = (item + (0, None, False))[:4]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = max(level, 0)
        bullet = "" if level == -1 else ("– " if level > 0 else "• ")
        p.text = f"{bullet}{text}" if level != -1 else text
        p.font.size = Pt(font_size - 2 * level if level else font_size)
        p.font.color.rgb = color or BLACK
        p.font.bold = bold
        p.space_after = Pt(8)
    return box


def add_image_fit(slide, path, left, top, max_w, max_h):
    im = Image.open(path)
    iw, ih = im.size
    ratio = iw / ih
    w, h = max_w, max_w / ratio
    if h > max_h:
        h = max_h
        w = max_h * ratio
    x = left + (max_w - w) / 2
    y = top + (max_h - h) / 2
    slide.shapes.add_picture(path, x, y, width=Emu(int(w)), height=Emu(int(h)))


def title_slide():
    slide = new_slide()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.5), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Evaluating Modular Representation & Recognizer Pipelines\nfor ASL Alphabet Recognition"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    p.text = "Computer Vision Capstone Project — 2025.2"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    auth = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11), Inches(1.2))
    tf = auth.text_frame
    tf.word_wrap = True
    lines = [
        "Tran Quang Hung - 20235502",
        "Do Dang Vu - 20235578",
        "Le Hoang Tung - 20235572",
    ]
    for i, l in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = l
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE

    inst = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11), Inches(0.5))
    p = inst.text_frame.paragraphs[0]
    p.text = "School of Information and Communication Technology — Hanoi University of Science and Technology"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def content_slide(title, section=None):
    slide = new_slide()
    add_header(slide, title, section)
    add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
title_slide()

# 1. Motivation
s = content_slide("Motivation", "Introduction")
add_image_fit(s, "photos/asl_examples_montage.png", Inches(7.3), Inches(1.15), Inches(5.5), Inches(3.1))
add_bullets(s, [
    "ASL is the primary communication medium for millions of Deaf and hard-of-hearing people.",
    "Real-time recognition with only a webcam — no specialized hardware.",
    ("Core question: does a “smarter” representation (hand-crop, landmarks, enhancement) actually "
     "improve recognition — or just add complexity and a new point of failure?", 0, RED, True),
], top=Inches(1.3), width=Inches(6.5), font_size=19)
box = s.shapes.add_textbox(Inches(7.3), Inches(4.3), Inches(5.5), Inches(1.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "8 of the 26 ASL alphabet classes (A–H), test_split samples"
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# 2. Objective & Contributions
s = content_slide("Objective & Contributions", "Introduction")
add_bullets(s, [
    ("Research question:", -1, NAVY, True),
    ("Do better representations (hand-crop, landmarks, enhancement) really improve A–Z recognition, "
     "or do they only add complexity without adding value?", 1, None, False),
], top=Inches(1.1), font_size=19)
cols = [
    ("1. Modular framework", "Representation and Recognizer are decoupled and freely swappable — 13 distinct pipelines evaluated under one protocol."),
    ("2. Multi-criteria evaluation", "Not just clean accuracy — macro-F1, robustness under 7–8 corruption types, hand-detection failure rate."),
    ("3. Deep visual analysis", "Grad-CAM inside the CNN, UMAP/t-SNE of the feature space, step-by-step qualitative grids."),
]
col_w = Inches(3.9)
for i, (head, body) in enumerate(cols):
    x = MARGIN + i * (col_w + Inches(0.15))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), col_w, Inches(3.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    box.line.color.rgb = BLUE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    p.text = head
    p.font.bold = True
    p.font.size = Pt(17)
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(14)
    p2.font.color.rgb = BLACK
    p2.space_before = Pt(8)

# 3. Dataset Overview
s = content_slide("Dataset Overview", "Dataset")
rows = [
    ("Dataset", "Role", "Scale"),
    ("ASL Alphabet (Kaggle) - train", "Train classifier / landmark MLP", "26 classes, ~78,000 img"),
    ("ASL Alphabet - test_split", "Clean + robustness evaluation", "26 classes, 2,600 img"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, MARGIN, Inches(1.15), Inches(5.6), Inches(1.3))
table = tbl_shape.table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.1)
table.columns[2].width = Inches(1.5)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12)
        cell.text_frame.word_wrap = True
        if r == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
add_bullets(s, [
    ("Studio photos, clean background, even lighting → an “easy” dataset.", 0, None, False),
    ("Many pipelines reach 99–100% clean accuracy — robustness benchmarks, not clean accuracy, "
     "actually differentiate pipelines.", 0, RED, True),
], top=Inches(2.6), width=Inches(5.6), font_size=16)
add_image_fit(s, "photos/embedding_landmarks_by_class.png", Inches(6.3), Inches(1.1), Inches(6.5), Inches(5.6))
box = s.shapes.add_textbox(Inches(6.3), Inches(6.75), Inches(6.5), Inches(0.5))
p = box.text_frame.paragraphs[0]
p.text = "UMAP / t-SNE of 42-d landmark features: the 26 classes form well-separated clusters."
p.font.size = Pt(12.5)
p.font.italic = True
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# 4. Methodology overview
s = content_slide("Methodology: A Two-Module Architecture", "Methodology")
box = s.shapes.add_textbox(MARGIN, Inches(1.3), Inches(12.3), Inches(0.6))
p = box.text_frame.paragraphs[0]
p.text = "RGB image  →  [ Representation Module ]  →  [ Recognizer Module ]  →  A–Z label"
p.font.size = Pt(20)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = NAVY

for i, (head, body, color) in enumerate([
    ("Module 1: Representation", "Turns the raw image into something the recognizer can use: a cropped image, an enhanced image, or a landmark feature vector.", BLUE),
    ("Module 2: Recognizer", "Takes the representation's output and returns one of 26 A–Z labels with top-k confidence.", RED),
]):
    x = MARGIN + i * Inches(6.2)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(6.0), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = head
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(15)
    p2.space_before = Pt(10)
add_bullets(s, [
    "13 pipelines = valid (Representation × Recognizer) combinations, evaluated under one protocol on the same test set.",
], top=Inches(4.7), font_size=18)

# 4.1 Representation module
s = content_slide("Module 1 — Representation (4 types)", "Methodology")
add_image_fit(s, "photos/grid_representation.png", MARGIN, Inches(1.05), Inches(12.3), Inches(5.0))
add_bullets(s, [
    ("Raw — resize only.   MediaPipe Crop — 21 keypoints → bbox → crop.   "
     "Landmarks — 42-d pose vector, no texture.   Enhancement — CLAHE/Gamma/Sharpen/Denoise.", 0, None, False),
    ("If MediaPipe finds no hand → representation returns None → recognizer outputs “Unknown” (always wrong).",
     0, RED, True),
], top=Inches(6.15), font_size=14)

# 4.2 Recognizer module
s = content_slide("Module 2 — Recognizer (5 types)", "Methodology")
rows = [
    ("Recognizer", "Model type", "Input"),
    ("ResNet18 (torchvision / resnet18_asl)", "CNN, fine-tuned on ASL", "224×224 image"),
    ("SigLIP ViT (HuggingFace)", "Vision Transformer, pretrained", "224×224 image"),
    ("Landmark MLP", "Small MLP", "42-d vector"),
    ("Landmark SVM", "Support Vector Machine", "42-d vector"),
    ("Landmark RF", "Random Forest", "42-d vector"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, MARGIN, Inches(1.3), Inches(12.3), Inches(2.6))
table = tbl_shape.table
table.columns[0].width = Inches(5.3)
table.columns[1].width = Inches(4.0)
table.columns[2].width = Inches(3.0)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(15)
        if r == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
add_bullets(s, [
    "ViT/SigLIP are pretrained on large-scale data — theoretical expectation: more accurate and more robust "
    "to noise than a small from-scratch-finetuned ResNet18.",
], top=Inches(4.2), font_size=18)

# 5. Evaluation metrics
s = content_slide("Evaluation Metrics: How Accuracy Is Actually Computed", "Evaluation Protocol")
add_bullets(s, [
    ("The problem: the codebase's “standard” accuracy is computed only over images where a hand WAS "
     "detected — it drops “Unknown” images from both numerator and denominator.", 0, RED, True),
], font_size=18)
box = s.shapes.add_textbox(MARGIN, Inches(1.9), Inches(12.3), Inches(1.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "clean_accuracy = correct / (images WITH a detected hand)        real_accuracy = correct / (ALL images, Unknown = wrong)"
p.font.size = Pt(15)
p.font.italic = True
p.font.color.rgb = NAVY
add_bullets(s, [
    "hand_detection_failure_rate = fraction of images where the representation returned None "
    "(only applies to MediaPipe-based pipelines).",
    "worst_corruption_accuracy = minimum accuracy across 7–8 corruption types (low-light, overexposure, "
    "motion-blur, gaussian-noise, rotation, scale-shift, partial-occlusion, crop-shift).",
    ("Consequence: a pipeline can have very high clean_accuracy but a much lower real accuracy — "
     "see the case study in the Results section.", 0, RED, True),
], top=Inches(3.1), font_size=18)

# 6. Quantitative results table
s = content_slide("Quantitative Results — 13 Pipelines", "Results")
rows = [
    ("Pipeline", "Clean Acc", "Macro-F1", "Worst-case", "Hand-fail", "FPS"),
    ("raw_siglip", "1.0000", "1.0000", "0.8912", "0.0%", "6.3"),
    ("enhancement_gamma_vit", "1.0000", "1.0000", "0.8927", "0.0%", "22.9"),
    ("enhancement_clahe_vit", "0.9996", "0.9996", "0.8629", "0.0%", "9.3"),
    ("raw_resnet18", "0.9985", "0.9985", "0.9202", "0.0%", "222.3"),
    ("enhancement_gamma_resnet18", "0.9985", "0.9985", "0.9135", "0.0%", "21.8"),
    ("enhancement_sharpen_resnet18", "0.9965", "0.9965", "0.1010", "0.0%", "244.6"),
    ("enhancement_clahe_resnet18", "0.9954", "0.9954", "0.7356", "0.0%", "16.7"),
    ("enhancement_denoise_resnet18", "0.9954", "0.9954", "0.7012", "0.0%", "7.9"),
    ("mediapipe_crop_vit", "0.9851", "0.9707", "0.0000", "17.15%", "11.2"),
    ("mediapipe_landmarks_rf", "0.9741", "0.9552", "0.0000", "17.15%", "16.6"),
    ("mediapipe_landmarks_svm", "0.9706", "0.9723", "0.0000", "17.15%", "19.8"),
    ("mediapipe_landmarks_mlp", "0.9705", "0.9556", "0.0000", "17.15%", "35.0"),
    ("mediapipe_crop_resnet18", "0.8859", "0.8627", "0.0000", "17.15%", "18.9"),
]
tbl_shape = s.shapes.add_table(len(rows), 6, MARGIN, Inches(1.1), Inches(12.3), Inches(5.0))
table = tbl_shape.table
table.columns[0].width = Inches(4.0)
for c in range(1, 6):
    table.columns[c].width = Inches(1.66)
highlight_rows = {9, 13}  # mediapipe_crop_vit, mediapipe_crop_resnet18 (1-indexed including header => actual idx)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12.5)
        if r == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            if c == 3 and val == "0.0000":
                para.font.color.rgb = RED
                para.font.bold = True
            if r in highlight_rows:
                para.font.bold = True

footnote = s.shapes.add_textbox(MARGIN, Inches(6.25), Inches(12.3), Inches(0.5))
p = footnote.text_frame.paragraphs[0]
p.text = "Every MediaPipe-based pipeline collapses to 0% accuracy on the worst corruption — even pipelines with high clean accuracy."
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY

# 6.1-6.3 charts
chart_slides = [
    ("Accuracy Comparison Across Pipelines", "photos/pipeline_accuracy_comparison.png", None),
    ("Latency / FPS Trade-off", "photos/pipeline_latency_comparison.png",
     "Landmark-based (MLP/SVM/RF) and ResNet18-on-raw reach the highest FPS — suitable for real-time use."),
    ("Robustness Under Corruption", "photos/pipeline_robustness_comparison.png", None),
]
for title, img, note in chart_slides:
    s = content_slide(title, "Results")
    add_image_fit(s, img, MARGIN, Inches(1.15), Inches(12.3), Inches(5.4) if note else Inches(5.9))
    if note:
        box = s.shapes.add_textbox(MARGIN, Inches(6.7), Inches(12.3), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = GRAY

# 6.4 Qualitative representation grid
s = content_slide("Qualitative — Representation Module", "Results")
add_image_fit(s, "photos/grid_representation.png", MARGIN, Inches(1.1), Inches(12.3), Inches(5.5))
box = s.shapes.add_textbox(MARGIN, Inches(6.7), Inches(12.3), Inches(0.5))
p = box.text_frame.paragraphs[0]
p.text = ("Same 5 sample images, 7 transformations. The MediaPipe Crop/Landmarks columns are black on 3/5 rows "
          "— the hand was not detected in those (dark) images.")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY

# 6.5 Qualitative recognizer grid
s = content_slide("Qualitative — Recognizer Module (All 13 Pipelines)", "Results")
add_image_fit(s, "photos/grid_recognizer.png", MARGIN, Inches(1.1), Inches(12.3), Inches(5.5))
box = s.shapes.add_textbox(MARGIN, Inches(6.7), Inches(12.3), Inches(0.5))
p = box.text_frame.paragraphs[0]
p.text = ("Green = correct, red = wrong. Every MediaPipe-based pipeline is wrong on the exact same 3 dark "
          "images — the failure propagates from the shared representation stage.")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY

# 6.6 Case study
s = content_slide("Case Study — “High Accuracy, Yet It Still Collapses”", "Results")
add_image_fit(s, "photos/compare_A909_hand_not_detected.png", MARGIN, Inches(1.1), Inches(6.8), Inches(5.9))
box = s.shapes.add_textbox(Inches(7.6), Inches(1.2), Inches(5.2), Inches(5.6))
tf = box.text_frame
tf.word_wrap = True
lines = [
    ("mediapipe_crop_vit: Clean Accuracy = 98.5%", 17, NAVY, True),
    ("But the real accuracy (counting Unknown as wrong):", 16, BLACK, False),
    ("2110 / 2600 = 81.15%", 22, RED, True),
    ("17 percentage points lower than the headline number — because 17.15% of images where MediaPipe "
     "missed the hand are excluded from the standard accuracy.", 14, GRAY, False),
    ("raw_resnet18 (no detection step) still correctly predicts “A” with 100% confidence on this exact image.", 14, GREEN, True),
]
for i, (text, size, color, bold) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(14)

# 6.7 Grad-CAM
s = content_slide("Looking Inside the CNN — Grad-CAM", "Results")
add_image_fit(s, "photos/cnn_internals_raw_resnet18_P979.png", MARGIN, Inches(1.2), Inches(12.3), Inches(2.6))
lbl = s.shapes.add_textbox(MARGIN, Inches(1.0), Inches(4), Inches(0.3))
lbl.text_frame.paragraphs[0].text = "raw_resnet18"
lbl.text_frame.paragraphs[0].font.size = Pt(13)
lbl.text_frame.paragraphs[0].font.bold = True
add_image_fit(s, "photos/cnn_internals_mediapipe_crop_resnet18_Z911.png", MARGIN, Inches(4.1), Inches(12.3), Inches(2.0))
lbl2 = s.shapes.add_textbox(MARGIN, Inches(3.9), Inches(4), Inches(0.3))
lbl2.text_frame.paragraphs[0].text = "mediapipe_crop_resnet18"
lbl2.text_frame.paragraphs[0].font.size = Pt(13)
lbl2.text_frame.paragraphs[0].font.bold = True
box = s.shapes.add_textbox(MARGIN, Inches(6.55), Inches(12.3), Inches(0.6))
p = box.text_frame.paragraphs[0]
p.text = ("Activation maps shrink from conv1 to layer4 (224→7), channels grow 64→512. "
          "Grad-CAM correctly highlights the hand region driving the prediction.")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY
box.text_frame.word_wrap = True

# 6.8 Embedding by class (CNN features — landmark version already shown on the Dataset slide)
s = content_slide("Feature Space — CNN Embeddings by Class", "Results")
add_image_fit(s, "photos/embedding_cnn_by_class.png", MARGIN, Inches(1.1), Inches(12.3), Inches(5.5))
box = s.shapes.add_textbox(MARGIN, Inches(6.7), Inches(12.3), Inches(0.5))
p = box.text_frame.paragraphs[0]
p.text = ("512-d ResNet18 avgpool embeddings, 30 images/class. Classes separate just as clearly in deep "
          "feature space as in raw landmark space (compare to the Dataset slide).")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY

# 6.8b More qualitative examples (additional comparisons across letters)
s = content_slide("More Qualitative Examples — B and L", "Results")
add_image_fit(s, "photos/compare_B909.png", MARGIN, Inches(1.1), Inches(6.0), Inches(5.6))
add_image_fit(s, "photos/compare_L909.png", Inches(6.7), Inches(1.1), Inches(6.0), Inches(5.6))
box = s.shapes.add_textbox(MARGIN, Inches(6.8), Inches(12.3), Inches(0.5))
p = box.text_frame.paragraphs[0]
p.text = "Easy cases: good lighting → every pipeline (raw, crop, landmark, enhancement, ViT) predicts correctly."
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY

# 6.9 Representation shift
s = content_slide("Feature Space — Does Representation Shift It?", "Results")
add_image_fit(s, "photos/embedding_representation_shift.png", Inches(1.8), Inches(1.1), Inches(9.7), Inches(5.5))
box = s.shapes.add_textbox(MARGIN, Inches(6.7), Inches(12.3), Inches(0.6))
p = box.text_frame.paragraphs[0]
p.text = ("Same images, 4 representations, same ResNet18 backbone. CLAHE/Gamma/Raw land in nearly the same "
          "tiny cluster (mild enhancement barely shifts the feature space) — MediaPipe Crop (orange) lands "
          "in a clearly different region: changing representation, not just enhancement, truly moves the feature space.")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = GRAY
box.text_frame.word_wrap = True

# 7. Conclusion
s = content_slide("Conclusion", "Conclusion")
cols = [
    ("Measure the right metric", "Standard “clean accuracy” can hide real-world failure — always also report accuracy counting Unknown as wrong."),
    ("Two-stage errors compound", "Adding a hand-detection step (MediaPipe) adds a new point of failure — robustness collapses to 0% even with a strong backbone."),
    ("Simpler can be more robust", "Raw/enhancement + ResNet18/ViT never drops below ~70% accuracy even under the worst corruption."),
]
col_w = Inches(3.9)
for i, (head, body) in enumerate(cols):
    x = MARGIN + i * (col_w + Inches(0.15))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.2), col_w, Inches(2.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    box.line.color.rgb = BLUE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = head
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(13)
    p2.space_before = Pt(8)

box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(3.8), Inches(12.3), Inches(1.1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFC, 0xEC, 0xEC)
box.line.color.rgb = RED
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(10)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "Limitations"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = RED
p2 = tf.add_paragraph()
p2.text = ("Some representations are unfinished stubs (YOLO-crop, segmentation-mask, fusion); evaluation uses "
           "a single dataset (ASL Alphabet, studio photos).")
p2.font.size = Pt(13)

box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(5.1), Inches(12.3), Inches(1.1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xEC, 0xF4, 0xFC)
box.line.color.rgb = BLUE
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(10)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "Future Work"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = BLUE
p2 = tf.add_paragraph()
p2.text = ("Finish YOLO-crop/segmentation/fusion representations; evaluate on real-world (non-studio) webcam "
           "footage; add a raw-image fallback when MediaPipe fails to detect a hand.")
p2.font.size = Pt(13)

# Work assignment
s = content_slide("Work Assignment", "Appendix")
rows = [
    ("Member", "Responsibilities"),
    ("Tran Quang Hung - 20235502", "YOLOv8 pipeline, dataset preparation, hand detection"),
    ("Do Dang Vu - 20235578", "Classifier (ResNet/MobileNet) & transforms, robustness evaluation, visualization"),
    ("Le Hoang Tung - 20235572", "Demo (OpenCV/Gradio), evaluation pipeline, report"),
]
tbl_shape = s.shapes.add_table(len(rows), 2, MARGIN, Inches(1.3), Inches(12.3), Inches(2.6))
table = tbl_shape.table
table.columns[0].width = Inches(4.0)
table.columns[1].width = Inches(8.3)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(15)
        cell.text_frame.word_wrap = True
        if r == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# References
s = content_slide("References", "Appendix")
refs = [
    "[1] C. Lugaresi et al., “MediaPipe: A Framework for Building Perception Pipelines,” arXiv:1906.08172, 2019.",
    "[2] K. He, X. Zhang, S. Ren, J. Sun, “Deep Residual Learning for Image Recognition,” CVPR, 2016.",
    "[3] X. Zhai et al., “Sigmoid Loss for Language Image Pre-Training (SigLIP),” ICCV, 2023.",
    "[4] Kaggle, “ASL Alphabet Dataset,” kaggle.com/datasets/grassknoted/asl-alphabet.",
    "[5] L. McInnes, J. Healy, J. Melville, “UMAP: Uniform Manifold Approximation and Projection,” arXiv:1802.03426, 2018.",
    "[6] L. van der Maaten, G. Hinton, “Visualizing Data using t-SNE,” JMLR, 2008.",
    "[7] R. Selvaraju et al., “Grad-CAM: Visual Explanations from Deep Networks via Gradient-based Localization,” ICCV, 2017.",
]
add_bullets(s, refs, font_size=15)

# Final / Contact (acts as conclusion recap + contact, stays on screen for Q&A)
s = content_slide("Thank You — Questions Welcome", "Contact")
add_bullets(s, [
    ("Key takeaway:", -1, NAVY, True),
    "Always report accuracy that counts detection failures as wrong, and validate robustness — "
    "not just clean-image accuracy — before claiming a pipeline is “better.”",
    ("Team:", -1, NAVY, True),
    "Computer Vision Capstone — Hanoi University of Science and Technology",
    "Github: https://github.com/",
], font_size=18)

out_path = os.path.join(os.path.dirname(__file__), "Slide_ASL_EN.pptx")
prs.save(out_path)
print(f"Saved -> {out_path} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
